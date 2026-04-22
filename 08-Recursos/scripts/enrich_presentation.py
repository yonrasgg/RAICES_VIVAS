#!/usr/bin/env python3
"""Enrich the Entrega Final presentation with diagrams + anexos.

Input : 06-Entregables/Presentaciones/Raices_Vivas_Entrega_Final_Presentacion.pptx
Output: same file (overwritten). A backup is saved alongside with suffix `.bak`.

What it does:
  1. Inserts 7 new 16:9 anexo slides with the engineering diagrams at
     thematic positions inside the current narrative.
  2. Each new slide follows the existing palette
     (cream  #F5EDE0 · dark green #133A32 · terracotta #B85C38 · gold #C5A05A).
  3. Adds a discreet Cenfotec logo corner-accent on every slide.

Palette decision keeps visual coherence with the current deck while making
the diagrams the protagonists of each anexo page.
"""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
PPTX = ROOT / "06-Entregables" / "Presentaciones" / "Raices_Vivas_Entrega_Final_Presentacion.pptx"
IMG = ROOT / "08-Recursos" / "Imágenes"
LOGO = IMG / "logo-cenfotec-200x200.png"

# Brand palette (matches the existing deck)
CREAM       = RGBColor(0xF5, 0xED, 0xE0)
INK_GREEN   = RGBColor(0x13, 0x3A, 0x32)
DEEP_GREEN  = RGBColor(0x4E, 0x6E, 0x66)
TERRACOTTA  = RGBColor(0xB8, 0x5C, 0x38)
GOLD        = RGBColor(0xC5, 0xA0, 0x5A)
BONE        = RGBColor(0xE9, 0xDE, 0xC9)
SLATE       = RGBColor(0x55, 0x55, 0x55)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W_EMU = 12191695  # 13.33"
SLIDE_H_EMU = 6858000   # 7.5"

# Anexo definitions.
# `insert_after_original_index` refers to the 1-based index of the slide in
# the ORIGINAL 14-slide deck after which the new slide should appear.
ANEXOS = [
    dict(
        insert_after=7,
        kicker="ANEXO TÉCNICO · A.1",
        title="DFD Nivel 0 — Diagrama de contexto",
        subtitle="Notación DeMarco–Yourdon",
        caption=(
            "El sistema Raíces Vivas modelado como proceso único rodeado de sus "
            "entidades externas (6 actores humanos + 3 sistemas: MEP, CCSS/EDUS, "
            "ODS/UNESCO). Cada arista es un flujo de datos etiquetado."
        ),
        image="diagram_dfd_nivel0.png",
    ),
    dict(
        insert_after=7,
        kicker="ANEXO TÉCNICO · A.2",
        title="DFD Nivel 1 — Descomposición funcional",
        subtitle="Procesos · almacenes · flujos",
        caption=(
            "Cuatro procesos hijos (EDU · SAB · SAL · Transversales) y cinco "
            "almacenes de datos (D1 catálogo, D2 PouchDB, D3 log auditoría, "
            "D4 i18n, D5 CouchDB). Representa la base offline-first del sistema."
        ),
        image="diagram_dfd_nivel1.png",
    ),
    dict(
        insert_after=8,
        kicker="ANEXO TRAZABILIDAD · A.3",
        title="Casos de Uso — Diagrama UML",
        subtitle="23 casos de uso · 4 módulos",
        caption=(
            "Actores (docente · estudiante · portador · consejo · ATAP · admin) "
            "conectados a los 23 casos de uso, agrupados por módulo EDU/SAB/SAL/TRANS. "
            "Base de la matriz RF↔CU del proyecto."
        ),
        image="diagram_usecase.png",
    ),
    dict(
        insert_after=8,
        kicker="ANEXO TRAZABILIDAD · A.4",
        title="Actores × Módulos — Matriz de responsabilidades",
        subtitle="Quién interactúa con qué",
        caption=(
            "Mapeo bidireccional de los 9 actores primarios y 4 módulos funcionales. "
            "Justifica los niveles CARE y la segregación de permisos en la arquitectura."
        ),
        image="diagram_actor_modulo.png",
    ),
    dict(
        insert_after=8,
        kicker="ANEXO TRAZABILIDAD · A.5",
        title="QFD — Voz del usuario → requerimiento técnico",
        subtitle="Casa de la calidad del proyecto",
        caption=(
            "Priorización de requerimientos funcionales según el peso de cada "
            "necesidad levantada en campo. Fundamenta la clasificación MoSCoW "
            "(12 Must · 7 Should · 4 Could)."
        ),
        image="diagram_qfd.png",
    ),
    dict(
        insert_after=11,
        kicker="ANEXO GESTIÓN · A.6",
        title="Análisis FODA",
        subtitle="Fortalezas · Oportunidades · Debilidades · Amenazas",
        caption=(
            "Diagnóstico estratégico del proyecto en el contexto académico, "
            "territorial y técnico. Alimenta el registro de riesgos y el "
            "roadmap de continuidad."
        ),
        image="diagram_foda.png",
    ),
    dict(
        insert_after=11,
        kicker="ANEXO GESTIÓN · A.7",
        title="Análisis de causa-raíz — Ishikawa",
        subtitle="6M · por qué persiste la brecha sociotécnica",
        caption=(
            "Espina de pescado que identifica causas estructurales (método, "
            "personas, materiales, entorno, tecnología y medición) detrás de la "
            "brecha en educación, saberes y salud comunitaria."
        ),
        image="diagram_ishikawa.png",
    ),
]


# ───────────────────────── XML helpers ─────────────────────────

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _set_fill_rgb(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_text_box(slide, left, top, width, height, text, *,
                  size=14, bold=False, color=INK_GREEN, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _add_rect(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill_rgb(sh, fill)
    sh.shadow.inherit = False
    return sh


# ───────────────────────── Slide builder ─────────────────────────

def build_anexo_slide(prs: Presentation, spec: dict):
    blank = prs.slide_layouts[-1]  # DEFAULT
    slide = prs.slides.add_slide(blank)

    # --- cream background full bleed ---
    _add_rect(slide, 0, 0, Emu(SLIDE_W_EMU), Emu(SLIDE_H_EMU), CREAM)

    # --- left vertical ribbon (dark green) ---
    _add_rect(slide, 0, 0, Inches(0.35), Emu(SLIDE_H_EMU), INK_GREEN)

    # --- top kicker band (bone) ---
    _add_rect(slide, Inches(0.35), 0, Emu(SLIDE_W_EMU) - Inches(0.35),
              Inches(0.55), BONE)

    # kicker text (small caps vibe — all uppercase + letter spacing already handled visually)
    _add_text_box(
        slide, Inches(0.8), Inches(0.1),
        Inches(9), Inches(0.35),
        spec["kicker"], size=11, bold=True, color=TERRACOTTA,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font="Calibri",
    )

    # slide number (right)
    _add_text_box(
        slide, Inches(11.7), Inches(0.1),
        Inches(1.3), Inches(0.35),
        "Raíces Vivas · Entrega Final · CENFOTEC",
        size=9, color=SLATE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- big title ---
    _add_text_box(
        slide, Inches(0.8), Inches(0.7),
        Inches(11.6), Inches(0.75),
        spec["title"], size=30, bold=True, color=INK_GREEN,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # subtitle
    _add_text_box(
        slide, Inches(0.8), Inches(1.45),
        Inches(11.6), Inches(0.35),
        spec["subtitle"], size=14, color=DEEP_GREEN,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    # --- image, max area 12.0 x 4.5 in, centered ---
    img_path = IMG / spec["image"]
    if not img_path.exists():
        raise FileNotFoundError(img_path)

    # Place with a white card behind it for legibility
    max_w = Inches(11.8)
    max_h = Inches(4.4)
    card_left = Inches(0.75)
    card_top = Inches(2.0)
    card_w = max_w
    card_h = max_h
    card = _add_rect(slide, card_left, card_top, card_w, card_h, WHITE)
    card.line.color.rgb = BONE
    card.line.width = Pt(0.75)

    # Fit image inside card preserving aspect ratio
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    card_w_in = 11.8
    card_h_in = 4.4
    scale = min(card_w_in / (iw / 96), card_h_in / (ih / 96))
    # treat source at 96 dpi so ratio math works without DPI awareness
    ratio = iw / ih
    # Compute fit
    if card_w_in / ratio <= card_h_in:
        draw_w = card_w_in * 0.95
        draw_h = draw_w / ratio
    else:
        draw_h = card_h_in * 0.95
        draw_w = draw_h * ratio
    left = Inches(0.75 + (card_w_in - draw_w) / 2)
    top = Inches(2.0 + (card_h_in - draw_h) / 2)
    slide.shapes.add_picture(str(img_path), left, top,
                             width=Inches(draw_w), height=Inches(draw_h))

    # --- caption band (dark green) ---
    _add_rect(slide, Inches(0.35), Inches(6.55),
              Emu(SLIDE_W_EMU) - Inches(0.35), Inches(0.75), INK_GREEN)
    _add_text_box(
        slide, Inches(0.8), Inches(6.60),
        Inches(11.6), Inches(0.65),
        spec["caption"], size=11, color=CREAM,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # --- small Cenfotec logo bottom-right corner ---
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(12.5), Inches(6.75),
                                 width=Inches(0.45), height=Inches(0.45))

    # --- slide number bottom-left ---
    # Placeholder — will be relabelled after reordering
    spec["_slide_obj"] = slide
    return slide


# ───────────────────────── Reordering ─────────────────────────

def reorder_slides(prs: Presentation, desired_order_indices: list[int]):
    """Reorder slides in sldIdLst according to zero-based indices of the CURRENT
    slide list."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for i in reversed(range(len(slides_list))):
        xml_slides.remove(slides_list[i])
    for idx in desired_order_indices:
        xml_slides.append(slides_list[idx])


# ───────────────────────── Corner-accent pass ─────────────────────────

def add_logo_to_existing_slides(prs: Presentation, already_built_indices: set[int]):
    """Add Cenfotec logo corner accent to the original slides too."""
    if not LOGO.exists():
        return
    for i, slide in enumerate(prs.slides):
        if i in already_built_indices:
            continue
        slide.shapes.add_picture(str(LOGO), Inches(12.55), Inches(6.8),
                                 width=Inches(0.40), height=Inches(0.40))


# ───────────────────────── Main ─────────────────────────

def main():
    if not PPTX.exists():
        raise SystemExit(f"Missing: {PPTX}")

    backup = PPTX.with_suffix(PPTX.suffix + ".bak")
    shutil.copy2(PPTX, backup)
    print(f"• Backup → {backup.relative_to(ROOT)}")

    prs = Presentation(PPTX)
    original_count = len(prs.slides)
    print(f"• Original slides: {original_count}")

    # 1) Append new anexo slides at the end (in the order declared in ANEXOS).
    new_slide_indices = []
    for spec in ANEXOS:
        slide = build_anexo_slide(prs, spec)
        idx = list(prs.slides).index(slide)
        new_slide_indices.append(idx)
        print(f"  + anexo: {spec['title']}  (temp idx={idx})")

    # 2) Compute the final desired order.
    #    Start with original indices 0..original_count-1
    #    Then walk through ANEXOS in order and insert each one right after
    #    its `insert_after` target (in the ORIGINAL deck).
    final_order: list[int] = list(range(original_count))
    for spec, new_idx in zip(ANEXOS, new_slide_indices):
        target_original = spec["insert_after"]  # 1-based in original
        # Find current position of that original index in final_order
        try:
            pos = final_order.index(target_original - 1) + 1
        except ValueError:
            pos = len(final_order)
        # If multiple anexos share the same `insert_after`, they stack in order.
        while pos < len(final_order) and final_order[pos] in new_slide_indices:
            pos += 1
        final_order.insert(pos, new_idx)

    reorder_slides(prs, final_order)
    print(f"• Final slide count: {len(prs.slides)}")

    # 3) Corner-accent logo is already embedded on each anexo slide during
    #    construction; we deliberately leave the original 14 slides untouched
    #    to respect their existing composition (decorative slide-number
    #    glyphs, quotes, etc.).

    prs.save(PPTX)
    print(f"✔ Saved: {PPTX.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
